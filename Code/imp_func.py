from tokenize import group
import pandas as pd
import os
import snowflake.connector
from pptx.chart.data import ChartData
from pptx.util import Pt,Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
import calendar
from datetime import datetime, timedelta,date
from dateutil import relativedelta
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns


def calculate_chart(chart_object,df,ver,col_list):
    chart_data = ChartData()
    chart_data.categories = df[ver].apply(lambda x: x.date()).values.tolist()
    for col in col_list:
        if col in df.columns:
            chart_data.add_series(str(col), df[col].values.tolist())
        else:
            continue
    chart_object.chart.replace_data(chart_data)
    chart_object.chart.has_legend = True

def update_chart(chart_object,df,ver,col_list):
    chart_data = ChartData()
    chart_data.categories = df[ver].apply(lambda x: str(x)).values.tolist()
    for col in col_list:
        if col in df.columns:
            chart_data.add_series(str(col), [x for x in df[col].values if str(x) != 'nan'])
        else:
           continue
    chart_object.chart.replace_data(chart_data)
    chart_object.chart.has_legend = True


def transform_month_number(df):
    df["MONTH"] = df["MONTH"].apply(lambda x: calendar.month_abbr[x])
    temp = df.fillna('').rename(columns={"MONTH": "Application Month"})
    return temp  


# define func to write cells
def write_cell(cell, text, font_name='Arial', font_size=Pt(11), font_bold=False):
    cell.text_frame.paragraphs[0].text = text
    cell.text_frame.paragraphs[0].font.name = font_name
    cell.text_frame.paragraphs[0].font.size = font_size
    cell.text_frame.paragraphs[0].font.bold = font_bold
    return

def write_table(table_instance, data_to_write, row_start, col_start,font_name=None, font_size=None, font_bold=None):
    col = col_start
    for row in data_to_write.values:
        for value in row:
            if font_name is not None:
                write_cell(table_instance.cell(row_start, col), value, font_name=font_name, font_size=font_size,
                           font_bold=font_bold)
            else:
                write_cell(table_instance.cell(row_start, col), value)
            col += 1
        row_start += 1
        col = col_start

    return


    